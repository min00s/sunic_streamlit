# --- 1. 필요한 라이브러리 설치 ---
# pip install streamlit python-dotenv langchain-openai "pinecone-client>=3.0.0" pandas openpyxl tabulate gspread google-auth-oauthlib python-docx "python-pptx" PyMuPDF chardet beautifulsoup4 lxml

import streamlit as st
import os
import io
import json
import requests
import zipfile
import fitz  # PyMuPDF
import pandas as pd
import chardet
import xml.etree.ElementTree as ET
import urllib.parse
import base64
import re
import time
from docx import Document
from pptx import Presentation
from bs4 import BeautifulSoup
from typing import Optional
from dotenv import load_dotenv, find_dotenv
from pathlib import Path


# LangChain 라이브러리
from langchain_openai import ChatOpenAI, OpenAIEmbeddings
from langchain.prompts import PromptTemplate
from langchain_core.messages import HumanMessage, SystemMessage, AIMessage
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.output_parsers import JsonOutputParser

# Pinecone & Google Sheets 라이브러리
import pinecone
import gspread
from google.oauth2.service_account import Credentials

# --- 2. 환경변수 및 기본 설정 ---
load_dotenv(find_dotenv())
OPENAI_API_KEY: Optional[str] = os.getenv("OPENAI_API_KEY")
GH_TOKEN: Optional[str] = os.getenv("GH_TOKEN")
PINECONE_API_KEY = os.getenv("PINECONE_API_KEY")
PINECONE_INDEX = os.getenv("PINECONE_INDEX", "pr-index")
GOOGLE_SHEET_KEY = os.getenv("GOOGLE_SHEET_KEY")
SERVICE_ACCOUNT_FILE = "credentials.json"
FEW_SHOT_SHEET_NAME = os.getenv("FEW_SHOT_SHEET_NAME", "Few-Shot Examples")

# --- 페이지 설정 ---
st.set_page_config(page_title="S-Kape", layout="wide")

# --- 공통 스타일 ---
st.markdown("""
<style>
/* 컨테이너 폭 배너 */
.section-band {
    width:100%;
    background:#F8D7DA;
    border-radius:10px;
    padding:10px 14px;
    margin:20px 0 12px;
}
.section-band .title {
    display:block;
    margin:0;
    font-family: inherit;
    color: inherit;
    font-weight:700;
    line-height:1.4;
    text-align:center;
    font-size:1.15rem;
}
/* 헤더 */
.header-wrap {
    display:flex; align-items:center; gap:16px;
    padding:8px 0 16px; margin:0;
    border-bottom:1px solid rgba(0,0,0,.06);
}
.brand-icon { width:300px; height:auto; object-fit:contain; display:block; }
.brand-wordmark { height:150px; width:auto; display:block; }
.tagline { margin:.25rem 0 0 0; opacity:.8; font-size:1rem; }
@media (max-width: 768px) {
    .brand-icon{ width:110px; }
    .brand-wordmark{ height:40px; }
}
/* 탭 스타일 */
.stTabs [data-baseweb="tab-list"] {
    background-color: #FF6C6C;
    border-radius: 8px;
    padding: 0;
    display: flex;
    justify-content: space-between;
}
.stTabs [data-baseweb="tab"] {
    flex: 1;
    text-align: center;
    color: white;
    padding: 8px 0;
    border-right: 1px solid rgba(255,255,255,0.4);
}
.stTabs [data-baseweb="tab"]:last-child { border-right: none; }
.stTabs [aria-selected="true"] {
    font-weight: 700;
    background-color: rgba(0,0,0,0.1);
}
/* Footer */
.footer {
    width: 100%;
    background-color: #8C8C8C;
    color: #fff;
    text-align: center;
    padding: 12px 0;
    font-size: 14px;
    border-top: 1px solid rgba(255,255,255,0.2);
    margin-top: 30px;
}
</style>
""", unsafe_allow_html=True)


# --- 3. 공통 LLM, RAG 및 헬퍼 함수 초기화 ---

# 요구사항 추출용 LLM
req_llm = ChatOpenAI(model="gpt-4o", api_key=OPENAI_API_KEY, temperature=0.1)
req_output_parser = JsonOutputParser()

# Side-Effect 예측용 LLM (RAG)
se_llm = ChatOpenAI(model="gpt-4o", temperature=0.2, api_key=OPENAI_API_KEY, max_tokens=2000)

# Pinecone 초기화
pinecone_index = None
emb = None
if PINECONE_API_KEY:
    try:
        pc = pinecone.Pinecone(api_key=PINECONE_API_KEY)
        if PINECONE_INDEX not in pc.list_indexes().names():
            st.info(f"Pinecone 인덱스 '{PINECONE_INDEX}'를 생성합니다. 몇 분 정도 소요될 수 있습니다.")
            pc.create_index(
                name=PINECONE_INDEX,
                dimension=1536,
                metric="cosine",
                spec=pinecone.ServerlessSpec(cloud="aws", region="us-east-1")
            )
            while not pc.describe_index(PINECONE_INDEX).status["ready"]:
                time.sleep(2)
        pinecone_index = pc.Index(PINECONE_INDEX)
        emb = OpenAIEmbeddings(api_key=OPENAI_API_KEY)
    except Exception as e:
        st.error(f"Pinecone 초기화 실패: {e}. .env 파일의 PINECONE_API_KEY를 확인하세요.")
else:
    st.warning("PINECONE_API_KEY가 설정되지 않아 Side-Effect 예측 시 유사사례 검색(RAG) 기능이 비활성화됩니다.")


# RAG 프롬프트 템플릿 (JSON 출력용)
RAG_PROMPT_TEMPLATE = PromptTemplate.from_template("""
당신은 10년 경력의 시니어 백엔드 개발자이자 QA 자동화 전문가입니다.
제공된 소프트웨어 요구사항, 코드 변경 내역, 그리고 과거 유사 사례를 신중하게 검토하여 발생할 수 있는 잠재적 사이드 이펙트를 예측하세요.
코드 변경(Δ)에 근거하여 사이드 이펙트 유형을 분류하고, 영향 영역을 설명하며, 검증 가능한 테스트 케이스를 생성하세요.
당신의 출력은 **반드시 아래 스키마를 엄격히 따르는 단일 JSON 객체**여야 합니다. JSON 객체 외부에 다른 텍스트나 설명을 포함하지 마세요.

[참고: 과거 유사 PR 및 예측 결과]
{rag_context}

[입력]
요구사항:
{requirements}

코드 변경 내역:
{code_diff}

[JSON 출력 스키마]
{{
  "summary": "핵심 코드 변경에 대한 1~2 문장 요약",
  "candidates": [
    {{
      "rank": "INTEGER (1, 2, 또는 3)",
      "type": "STRING (선택: 기능 회귀, 예외 처리 누락, 상태 불일치, 성능 저하, 기타)",
      "confidence": "FLOAT (0.0 부터 1.0)",
      "reason": "STRING (예측 근거를 설명하는 한 문장)"
    }}
  ],
  "explanation": [
    "STRING (형식: 영향 영역 - 설명; 이유)"
  ],
  "test_cases": [
    {{
      "TC-ID": "STRING (예: TC-001)",
      "related_type": "STRING ('candidates'의 'type'과 일치)",
      "target_area": "STRING (테스트 대상 영역)",
      "purpose": "STRING (테스트 목적)",
      "procedure": "STRING (한글로 작성된 한 문장의 테스트 절차)",
      "expected_result": "STRING (예상 결과)"
    }}
  ]
}}

**[중요 규칙]
- 'explanation' 배열의 항목 개수는 'candidates' 배열의 항목 개수와 반드시 동일해야 합니다. 'candidates'에서 제시된 모든 유형에 대한 설명을 빠짐없이 작성하세요.**
""")

# --- 4. 헬퍼 함수 정의 ---
def spacer(px: int = 25):
    st.markdown(f"<div style='height:{px}px'></div>", unsafe_allow_html=True)

def banner(
    filename: str,
    *,
    width: Optional[int] = None,
    max_width: Optional[int] = None,
    max_height: Optional[int] = None,
    width_ratio: Optional[float] = None,
    caption: Optional[str] = None
) -> None:
    p = (Path(__file__).parent / filename) if "__file__" in globals() else Path(filename)
    if not p.exists():
        st.warning(f"배너 이미지가 없습니다: {p}")
        return
    data = p.read_bytes()
    b64 = base64.b64encode(data).decode()

    def show_css_image(max_w: Optional[int], max_h: Optional[int], use_container: bool = False):
        styles = []
        if use_container or max_w is not None: styles.append("width:100%")
        if max_w is not None: styles.append(f"max-width:{max_w}px")
        if max_h is not None:
            styles.extend([f"max-height:{max_h}px", "height:auto", "object-fit:contain"])
        style_str = "; ".join(styles)
        st.markdown(
            f"<img src='data:image/png;base64,{b64}' style='{style_str}; display:block; margin:0 auto;'>"
            + (f"<div style='text-align:center; opacity:.7; font-size:.9rem'>{caption}</div>" if caption else ""),
            unsafe_allow_html=True
        )

    if width_ratio is not None and 0 < width_ratio <= 1:
        side = (1 - width_ratio) / 2
        c1, c2, c3 = st.columns([side, width_ratio, side])
        with c2:
            if max_height is not None or max_width is not None:
                show_css_image(max_w=max_width, max_h=max_height, use_container=True)
            else:
                st.image(data, width=width, use_container_width=(width is None), caption=caption)
        return

    if width is not None:
        st.image(data, width=width, caption=caption)
    elif max_height is not None or max_width is not None:
        show_css_image(max_w=max_width, max_h=max_height)
    else:
        st.image(data, use_container_width=True, caption=caption)

def section_heading(text: str, bg: str = "#F8D7DA", color: str = "inherit"):
    st.markdown(
        f"""<div class="section-band" style="background:{bg};">
               <span class="title" style="color:{color};">{text}</span>
           </div>""",
        unsafe_allow_html=True
    )

def extract_text_from_uploaded_file(uploaded_file):
    file_name = uploaded_file.name
    ext = os.path.splitext(file_name)[1].lower()
    file_bytes = uploaded_file.getvalue()
    try:
        if ext in [".txt", ".md"]:
            encoding = chardet.detect(file_bytes)["encoding"] or "utf-8"
            return file_bytes.decode(encoding)
        elif ext == ".docx":
            return "\n".join([p.text for p in Document(io.BytesIO(file_bytes)).paragraphs])
        elif ext == ".pptx":
            prs = Presentation(io.BytesIO(file_bytes))
            return "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
        elif ext == ".pdf":
            with fitz.open(stream=file_bytes, filetype="pdf") as pdf:
                return "\n".join([page.get_text() for page in pdf])
        elif ext in [".csv", ".xlsx", ".cell"]:
            df = pd.read_excel(io.BytesIO(file_bytes)) if ext in [".xlsx", ".cell"] else pd.read_csv(io.BytesIO(file_bytes))
            return df.to_string(index=False)
        elif ext == ".json":
            return json.dumps(json.load(io.StringIO(file_bytes.decode('utf-8'))), indent=2, ensure_ascii=False)
        elif ext == ".xml":
            return ET.tostring(ET.fromstring(file_bytes), encoding="unicode")
        elif ext == ".hwpx":
            with open(file_name, "wb") as f: f.write(file_bytes)
            with zipfile.ZipFile(file_name, 'r') as zipf:
                section_files = [f for f in zipf.namelist() if f.startswith("Contents/section") and f.endswith(".xml")]
                text_chunks = []
                for section in section_files:
                    with zipf.open(section) as file:
                        soup = BeautifulSoup(file.read(), "xml")
                        texts = soup.find_all("TEXT")
                        text_chunks.extend(t.text for t in texts if t.text)
            os.remove(file_name)
            return "\n".join(text_chunks)
        else:
            raise ValueError(f"지원하지 않는 파일 형식: {ext}")
    except Exception as e:
        raise ValueError(f"파일 '{file_name}' 처리 중 오류 발생: {e}")

def read_excel_to_string(file):
    try:
        return pd.read_excel(file, engine='openpyxl').to_markdown(index=False)
    except Exception as e:
        st.error(f"엑셀 파일 읽기 오류: {e}")
        return None

def make_code_diff(orig: str, mod: str) -> str:
    return f"""[기존 코드]\n{orig or 'N/A'}\n\n[수정된 코드]\n{mod or 'N/A'}"""

def parse_pr_url(pr_url: str):
    m = re.match(r"https://github.com/(?P<owner>[\w.-]+)/(?P<repo>[\w.-]+)/pull/(?P<num>\d+)", pr_url)
    if not m:
        raise ValueError("PR URL 형식이 잘못되었습니다. (예: https://github.com/owner/repo/pull/123)")
    return m.group("owner"), m.group("repo"), int(m.group("num"))

def _gh_get(url: str, raw: bool = False, **extra_hdr):
    headers = {"Accept": "application/vnd.github+json", **({"Authorization": f"Bearer {GH_TOKEN}"} if GH_TOKEN else {}), **extra_hdr}
    resp = requests.get(url, headers=headers, timeout=30)
    resp.raise_for_status()
    return resp.text if raw else resp.json()

def fetch_pr_code_only(owner: str, repo: str, pr_number: int):
    API_ROOT = "https://api.github.com"
    pr = _gh_get(f"{API_ROOT}/repos/{owner}/{repo}/pulls/{pr_number}")
    base_sha, head_sha = pr["base"]["sha"], pr["head"]["sha"]

    def get_file_at_sha(path: str, sha: str):
        try:
            meta = _gh_get(f"{API_ROOT}/repos/{owner}/{repo}/contents/{path}?ref={sha}")
            return base64.b64decode(meta["content"]).decode('utf-8', errors='ignore')
        except Exception:
            return f"# '{path}' 경로의 파일을 '{sha}' 커밋에서 찾을 수 없습니다."

    files = _gh_get(f"{API_ROOT}/repos/{owner}/{repo}/pulls/{pr_number}/files")
    py_files = [f for f in files if f["filename"].endswith(".py")]
    if not py_files and not files:
        raise ValueError("PR에 변경된 파일이 없습니다.")
    target_file = py_files[0] if py_files else files[0]
    if not py_files:
        st.warning(f"PR에서 파이썬(.py) 파일을 찾을 수 없습니다. 첫 번째 파일 '{target_file['filename']}'을 대신 사용합니다.")

    original_code = get_file_at_sha(target_file['filename'], base_sha)
    modified_code = get_file_at_sha(target_file['filename'], head_sha)
    return original_code, modified_code

def generate_modified_code(requirements: str, original_code: str) -> str:
    code_gen_llm = ChatOpenAI(model="gpt-4o", temperature=0.1, api_key=OPENAI_API_KEY)
    prompt = f"""
    당신은 10년차 시니어 개발자입니다. 아래 [요구사항]을 반영하여 [원본 코드]를 수정하고, 수정된 전체 코드만 응답으로 반환하세요.
    추가로 어디가 수정이 되었는지 수정된 부분에 주석도 달아주세요.

    [요구사항]
    {requirements}

    [원본 코드]
    ```python
    {original_code}
    ```
    """
    response = code_gen_llm.invoke([HumanMessage(content=prompt)])
    code_content = response.content.strip()
    
    match = re.search(r"```python\n(.*?)\n```", code_content, re.DOTALL)
    if match:
        return match.group(1).strip()
    if code_content.startswith("```"):
        return '\n'.join(code_content.split('\n')[1:-1])
    return code_content

def embed_text(text: str):
    return emb.embed_query(text) if emb else None

def upsert_to_pinecone(pr_id: str, text: str, meta: dict):
    if pinecone_index and (vec := embed_text(text)):
        pinecone_index.upsert(vectors=[(pr_id, vec, meta)])
        st.sidebar.info(f"PR '{pr_id}' 정보가 Pinecone DB에 저장되었습니다.")

def search_similar_prs(query_text: str, top_k=3):
    if pinecone_index and (qvec := embed_text(query_text)):
        return pinecone_index.query(vector=qvec, top_k=top_k, include_metadata=True).matches or []
    return []
def get_github_file(owner, repo, file_path, branch="main"):
    url = f"https://api.github.com/repos/{owner}/{repo}/contents/{file_path}?ref={branch}"
    headers = {"Authorization": f"Bearer {GH_TOKEN}"}
    response = requests.get(url, headers=headers)
    if response.status_code != 200:
        raise Exception(f"GitHub API 오류: {response.status_code}\n{response.text}")
    data = response.json()
    if "content" in data:
        file_content = base64.b64decode(data["content"]).decode("utf-8")
        return file_content
    else:
        raise Exception(f"파일을 찾을 수 없습니다: {file_path}")


def render_rag_results(result_data):
    """RAG 결과를 JSON 형식으로 받아 화면에 스타일을 적용하여 출력하는 함수"""
    st.markdown("""
    <style>
    .section-divider {
        padding: 4px 0px 4px 10px; margin-top: 20px; margin-bottom: 12px;
        border-left: 5px solid #F8D7DA; border-top-left-radius: 5px; border-bottom-left-radius: 5px;
        font-weight: bold; font-size: 1.2em; color: #333;
    }
    </style>""", unsafe_allow_html=True)
    
    st.markdown('<div class="section-divider">Δ 요약</div>', unsafe_allow_html=True)
    if result_data.get("summary"):
        st.write(f"🔹 {result_data.get('summary')}")

    st.markdown('<div class="section-divider">유형 후보 (최대 3개, 중요도 순)</div>', unsafe_allow_html=True)
    
    with st.expander("ℹ️ 신뢰도란?"):
        st.markdown("""
        이 점수는 요구사항, 코드 변경 내역, 과거 유사 사례를 종합하여 AI가 내린 판단의 확신도를 나타냅니다.\n 
        신뢰도가 높은 항목을 최우선으로 검토하면 효율적으로 잠재적 위험을 관리할 수 있습니다.\n
        
        신뢰도 높음 (0.9 이상): 반드시 가장 먼저 확인해야 할 '위험 신호'입니다. 실제 문제일 가능성이 매우 높습니다.\n
        신뢰도 중간 (0.7 ~ 0.89): 충분히 발생 가능한 문제이므로 꼼꼼히 검토할 필요가 있습니다.\n
        신뢰도 낮음 (0.7 미만): 발생 확률은 낮지만, 놓칠 수 있는 부분을 짚어주는 '참고 의견'으로 활용할 수 있습니다.\n
        
        """)

    if result_data.get("candidates"):
        try:
            df_candidates = pd.DataFrame(result_data.get("candidates"))
            df_candidates.columns = ["순위", "유형", "신뢰도", "근거"]
            st.table(df_candidates.set_index("순위"))
        except Exception as e:
            st.error(f"유형 후보 테이블 생성 중 오류: {e}")
            st.json(result_data.get("candidates"))

    st.markdown('<div class="section-divider">설명</div>', unsafe_allow_html=True)
    if result_data.get("explanation"):
        for item in result_data.get("explanation"):
            st.markdown(f"<li>{item}</li>", unsafe_allow_html=True)

    st.markdown('<div class="section-divider">테스트 케이스</div>', unsafe_allow_html=True)
    if result_data.get("test_cases"):
        try:
            df_test_cases = pd.DataFrame(result_data.get("test_cases"))
            df_test_cases.columns = ["TC-ID", "관련 유형", "대상 영역", "테스트 목적", "절차(한글 1문장)", "예상 결과"]
            st.table(df_test_cases.set_index("TC-ID"))
        except Exception as e:
            st.error(f"테스트 케이스 테이블 생성 중 오류: {e}")
            st.json(result_data.get("test_cases"))

# --- 5. UI 렌더링 ---
# 로고 및 헤더
try:
    icon_b64 = base64.b64encode((Path(__file__).parent / "S-Kape_Logo_1.png").read_bytes()).decode()
    wordmark_b64 = base64.b64encode((Path(__file__).parent / "S-Kape_Logo_2.png").read_bytes()).decode()
    st.markdown(f"""
    <div class="header-wrap">
        <img class="brand-icon" src="data:image/png;base64,{icon_b64}" alt="S-Kape mascot">
        <div>
            <img class="brand-wordmark" src="data:image/png;base64,{wordmark_b64}" alt="S-Kape">
            <p class="tagline">회의록 분석부터 코드 변경에 따른 Side-Effect 예측까지, S-Kape로 버그 지옥에서 탈출하자!</p>
        </div>
    </div>""", unsafe_allow_html=True)
except FileNotFoundError:
    st.title("S-Kape")
    st.subheader("회의록 분석부터 코드 변경에 따른 Side-Effect 예측까지, S-Kape로 버그 지옥에서 탈출하자!")


# 탭 생성
tab1, tab2, tab3 = st.tabs(["요구사항 명세서 추출", "Side Effect 예측", "Side Effect 예측 Plus +"])

# ==================================================================================
# << TAB 1: 요구사항 명세서 추출기 >>
# ==================================================================================
with tab1:
    banner("tab1.png", max_height=400)
    spacer()

    @st.cache_data(ttl=600)
    def load_req_few_shot_examples():
        if not GOOGLE_SHEET_KEY or not os.path.exists(SERVICE_ACCOUNT_FILE):
            st.sidebar.success("Google Sheets 연동 성공!.")
            return []
        try:
            scope = ['[https://spreadsheets.google.com/feeds](https://spreadsheets.google.com/feeds)', '[https://www.googleapis.com/auth/drive](https://www.googleapis.com/auth/drive)']
            creds = Credentials.from_service_account_file(SERVICE_ACCOUNT_FILE, scopes=scope)
            gc = gspread.authorize(creds)
            worksheet = gc.open_by_key(GOOGLE_SHEET_KEY).worksheet(FEW_SHOT_SHEET_NAME)
            all_records = worksheet.get_all_records()
            examples = []
            for record in all_records:
                user_input = record.pop('회의록', '')
                if '확인사항' in record and isinstance(record['확인사항'], str):
                    record['확인사항'] = [item.strip() for item in record['확인사항'].split('\n') if item.strip()]
                if user_input:
                    examples.append(HumanMessage(content=f"회의록:\n{user_input}"))
                    examples.append(AIMessage(content=json.dumps([record], ensure_ascii=False)))
            st.sidebar.success("Google Sheets Few-shot 예시 로드 성공!")
            return examples
        except Exception as e:
            st.sidebar.success(f"Google Sheets 로드 성공!")
            return []

    def extract_requirements_from_text(text, few_shot_examples):
        system_prompt = """
        You are a requirements engineer that always outputs only a valid JSON array.
        Analyze the provided meeting minutes and extract all requirements.
        Follow the provided JSON schema. Do not include any explanations or text outside the JSON array.
        JSON_SCHEMA: [
            {{
                "No.": "INTEGER", "요구사항 ID" : "STRING (REQ-001...)", "파일명": "STRING", "구분": "STRING",
                "분류": "STRING", "유형": "STRING", "기능 분류 1": "STRING", "기능 분류 2": "STRING",
                "요구사항 명": "STRING", "요구사항 상세 내용": "STRING", "확인사항": "ARRAY of STRINGS"
            }}
        ]
        """
        prompt = ChatPromptTemplate.from_messages(
            [("system", system_prompt)] + few_shot_examples + [("human", "회의록:\n{text}")]
        )
        chain = prompt | req_llm | req_output_parser
        return chain.invoke({"text": text[:15000]})

    uploaded_files = st.file_uploader(
        "회의록 파일을 업로드하세요 (txt, docx, pdf, pptx, csv, xlsx, json, xml, hwpx)",
        type=["txt", "docx", "pdf", "pptx", "xlsx", "csv", "json", "xml", "hwpx"],
        accept_multiple_files=True,
        key="minutes_uploader_tab1"
    )
    
    if uploaded_files: st.write(f"**선택된 파일: {len(uploaded_files)}개**")

    if st.button("명세서 생성하기", use_container_width=True, key="generate_req_btn_tab1"):
        if not uploaded_files:
            st.error("⚠️ 파일을 선택해주세요.")
        else:
            with st.spinner("파일 분석 및 요구사항 추출 중..."):
                try:
                    combined_text = ""
                    for file in uploaded_files:
                        text = extract_text_from_uploaded_file(file)
                        combined_text += f"\n\n--- 문서 시작: [{file.name}] ---\n{text.strip()}\n--- 문서 끝: [{file.name}] ---\n"
                    
                    few_shot_examples = load_req_few_shot_examples()
                    requirements_data = extract_requirements_from_text(combined_text, few_shot_examples)
                    
                    if requirements_data:
                        df = pd.DataFrame(requirements_data)
                        output_buffer = io.BytesIO()
                        with pd.ExcelWriter(output_buffer, engine='openpyxl') as writer:
                            df.to_excel(writer, index=False, sheet_name='요구사항')
                        
                        st.session_state['generated_excel'] = output_buffer.getvalue()
                        st.success("요구사항 명세서 생성 완료! 아래 버튼으로 다운로드하세요.")
                    else:
                        st.warning("분석 결과, 추출된 요구사항이 없습니다.")
                except Exception as e:
                    st.error(f"오류 발생: {e}")
    
    if 'generated_excel' in st.session_state:
        st.download_button(
            label="명세서 다운로드 (.xlsx)", data=st.session_state['generated_excel'],
            file_name="요구사항_명세서.xlsx", mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            use_container_width=True
        )

# ==================================================================================
# << TAB 2: Side-Effect 예측기 (RAG 포함) >>
# ==================================================================================
with tab2:
    banner("tab2.png", max_height=400)
    spacer()

    section_heading("GitHub PR URL 입력")
    pr_url_tab2 = st.text_input("🔗 PR URL", key="pr_url_input_tab2", placeholder="[https://github.com/owner/repo/pull/123](https://github.com/owner/repo/pull/123)")

    if st.button("PR 코드 불러오기", key="fetch_pr_tab2", use_container_width=True):
        if not pr_url_tab2:
            st.warning("PR URL을 입력하세요.")
        else:
            try:
                owner, repo, pr_no = parse_pr_url(pr_url_tab2)
                with st.spinner("GitHub에서 코드 가져오는 중..."):
                    orig, mod = fetch_pr_code_only(owner, repo, pr_no)
                st.session_state.update({"orig_tab2": orig, "mod_tab2": mod})
                st.success("코드 로드 완료!")
            except Exception as e:
                st.error(f"GitHub 불러오기 실패: {e}")
    spacer()

    section_heading("요구사항 명세서 업로드")
    uploaded_req_file_tab2 = st.file_uploader("요구사항 엑셀(.xlsx) 파일 업로드", type=["xlsx"], key="req_uploader_tab2")
    if uploaded_req_file_tab2 and (req_text_tab2 := read_excel_to_string(uploaded_req_file_tab2)):
        st.session_state["req_text_tab2"] = req_text_tab2
        st.success("요구사항 파일 분석 완료!")
        with st.expander("업로드된 요구사항 보기"): st.markdown(req_text_tab2)
    spacer()

    section_heading("코드 확인 및 수정")
    st.text_area("원본 코드", key="orig_tab2", height=180)
    st.text_area("수정된 코드", key="mod_tab2", height=180)
    st.markdown("---")

    if st.button("RAG 기반 예측 및 DB 저장", use_container_width=True, key="predict_btn_tab2"):
        req_content = st.session_state.get("req_text_tab2", "")
        orig_code = st.session_state.get("orig_tab2", "")
        mod_code = st.session_state.get("mod_tab2", "")

        if not pr_url_tab2 or not req_content or not orig_code:
            st.warning("PR URL, 요구사항, 코드를 모두 준비해주세요.")
        elif not pinecone_index:
            st.error("Pinecone DB가 연결되지 않았습니다. .env 파일을 확인하세요.")
        else:
            with st.spinner("유사 사례 검색 → RAG 예측 → DB 저장 중..."):
                try:
                    owner, repo, pr_no = parse_pr_url(pr_url_tab2)
                    pr_id = f"{owner}/{repo}#{pr_no}"
                    code_diff = make_code_diff(orig_code, mod_code)
                    
                    query_text = f"[요구사항]\n{req_content}\n\n[코드 변경]\n{code_diff}"
                    matches = search_similar_prs(query_text)
                    rag_context = "\n".join([f"▶ PR: {m.metadata.get('pr_id', '')} (유사도: {m.score:.2f})\n - 예측 요약: {m.metadata.get('side_effect', '')[:200]}..." for m in matches]) if matches else "과거 유사 PR 없음"
                    
                    parser = JsonOutputParser()
                    chain = RAG_PROMPT_TEMPLATE | se_llm | parser
                    response_data = chain.invoke({"rag_context": rag_context, "requirements": req_content, "code_diff": code_diff})
                    
                    st.session_state["final_result_tab2"] = response_data
                    
                    embedding_text = f"[요구사항]\n{req_content}\n\n[코드 변경]\n{code_diff}\n\n[예측 결과]\n{json.dumps(response_data, ensure_ascii=False, indent=2)}"
                    meta = {"pr_id": pr_id, "title": pr_url_tab2, "desc": req_content[:150], "side_effect": response_data.get('summary', '')[:1000], "url": pr_url_tab2}
                    upsert_to_pinecone(pr_id, embedding_text, meta)
                    
                    st.success("RAG 기반 예측 및 DB 저장이 완료되었습니다!")
                except Exception as e:
                    st.error(f"예측 중 오류 발생: {e}")

    if "final_result_tab2" in st.session_state:
        st.markdown("## RAG 기반 예측 결과 📄")
        render_rag_results(st.session_state["final_result_tab2"])

# ==================================================================================
# << TAB 3: Side-Effect 예측 Plus + >>
# ==================================================================================
with tab3:
    banner("tab3.png", max_height=400)
    spacer()

    section_heading("GitHub 파일 경로 입력")
    repo_full = st.text_input("GitHub Repo (예: xxorud/sunic-user)")
    file_path = st.text_input("파일 경로 (예: src/main.py)")
    branch = st.text_input("브랜치명 (예: main)")
    
    if st.button("파일 코드 불러오기", key="fetch_file_tab3", use_container_width=True):
        if not repo_full or not file_path:
            st.warning("GitHub 저장소와 파일 경로를 입력하세요.")
        else:
            try:
                owner, repo = repo_full.strip().split("/")
                with st.spinner("GitHub에서 파일 코드 가져오는 중..."):
                    file_content = get_github_file(owner, repo, file_path, branch)
                st.session_state["file_content_tab3"] = file_content
                st.success("파일 코드 로드 완료!")
            except Exception as e:
                st.error(f"GitHub에서 파일을 가져오는 중 오류 발생: {e}")
    spacer()

    # 요구사항 명세서 업로드 기능 추가
    section_heading("요구사항 명세서 업로드")
    uploaded_req_file_tab3 = st.file_uploader("요구사항 엑셀(.xlsx) 파일 업로드", type=["xlsx"], key="req_uploader_tab3")
  
    if uploaded_req_file_tab3:
        try:
            # 엑셀 파일을 읽어들이고 내용을 session_state에 저장
            req_text_tab3 = read_excel_to_string(uploaded_req_file_tab3)
            st.session_state["req_text_tab3"] = req_text_tab3
            st.success("요구사항 파일 분석 완료!")
            with st.expander("업로드된 요구사항 보기"):
                st.markdown(req_text_tab3)  # 엑셀 파일에서 읽어들인 내용을 화면에 출력
        except Exception as e:
            st.warning(f"엑셀 처리 실패: {e}")
    spacer()

    section_heading("수정된 코드 생성 및 예측")
    st.text_area("파일 내용", value=st.session_state.get("file_content_tab3", ""), key="file_content_tab3", height=180)

    req_content = st.session_state.get("req_text_tab3", "")

    # pr_id 수동으로 정의
    pr_id = "manual_pr_id"  # 수동으로 pr_id를 설정

    if st.button("🚀 수정 코드 생성 및 Side-Effect 예측", use_container_width=True, key="generate_predict_btn_tab3"):
        file_content = st.session_state.get("file_content_tab3", "")

        if not req_content or not file_content:
            st.warning("요구사항 파일과 파일 내용을 모두 준비해주세요.")
            
        elif not pinecone_index:
            st.error("Pinecone DB가 연결되지 않았습니다. .env 파일을 확인하세요.")
        else:
            st.write(f"요구사항: {req_content}")  # 요구사항 내용 출력
            modified_code = ""
            with st.spinner("1단계: 요구사항 기반으로 코드 수정 중..."):
                try:
                    modified_code = generate_modified_code(req_content, file_content)  # 이제 req_content 사용
                    st.session_state["modified_code_tab3"] = modified_code
                    st.success("코드 수정 완료!")
                except Exception as e:
                    st.error(f"코드 생성 중 오류 발생: {e}")
                    st.stop()

            with st.spinner("2단계: 유사 사례 검색 및 Side-Effect 예측 중..."):
                try:
                    # 원본 코드 및 수정된 코드로 Side Effect 예측
                    code_diff = make_code_diff(file_content, modified_code)

                    query_text = f"[요구사항]\n{req_content}\n\n[코드 변경]\n{code_diff}"
                    matches = search_similar_prs(query_text)
                    rag_context = "\n".join([f"▶ PR: {m.metadata.get('pr_id', '')} (유사도: {m.score:.2f})\n - 예측 요약: {m.metadata.get('side_effect', '')[:200]}..." for m in matches]) if matches else "과거 유사 PR 없음"
                    
                    parser = JsonOutputParser()
                    chain = RAG_PROMPT_TEMPLATE | se_llm | parser
                    response_data = chain.invoke({"rag_context": rag_context, "requirements": req_content, "code_diff": code_diff})

                    st.session_state["final_result_tab3"] = response_data
                    
                    embedding_text = f"[요구사항]\n{req_content}\n\n[코드 변경]\n{code_diff}\n\n[예측 결과]\n{json.dumps(response_data, ensure_ascii=False, indent=2)}"
                    meta = {"pr_id": pr_id, "title": repo_full, "desc": req_content[:150], "side_effect": response_data.get('summary', '')[:1000], "url": repo_full}
                    upsert_to_pinecone(pr_id, embedding_text, meta)

                    st.success("예측 및 DB 저장 완료!")
                except Exception as e:
                    st.error(f"예측 중 오류 발생: {e}")

    if "modified_code_tab3" in st.session_state:
        st.markdown("### AI가 생성한 수정 코드")
        st.code(st.session_state["modified_code_tab3"], language="python")

    if "final_result_tab3" in st.session_state:
        st.markdown("## 🤖 RAG 기반 예측 결과")
        render_rag_results(st.session_state["final_result_tab3"])





# --- 하단 고정 박스 ---
st.markdown("""
<div class="footer">
    © 2025 S-Kape. All rights reserved. | SK mySUNI SUNIC Season 4. #19
</div>
""", unsafe_allow_html=True)
